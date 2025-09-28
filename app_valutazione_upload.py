from flask import Flask, request, render_template_string, send_file, session
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from pdfminer.high_level import extract_text
import re, os

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 25 * 1024 * 1024  # 25MB
app.secret_key = os.environ.get("APP_SECRET_KEY", "dev-secret-key-change-me")

HTML = """
<!doctype html><html lang="it"><head><meta charset="utf-8">
<title>Valutazione Presentazioni · Upload</title>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<link href="https://fonts.googleapis.com/css2?family=Public+Sans:wght@300;400;600;800&display=swap" rel="stylesheet">
<style>
:root{--p:#F41163;--s:#F652B0;--m:#e9e9e9;--ink:#2b2b2b;--r:16px}
*{box-sizing:border-box}body{margin:0;font-family:'Public Sans',system-ui,Segoe UI,Roboto,Arial;color:var(--ink)}
header{padding:14px 16px;border-bottom:1px solid var(--m);background:#fff}
h1{font-size:18px;margin:0;font-weight:800}
.container{max-width:980px;margin:18px auto;padding:0 12px}
.card{border:1px solid var(--m);border-radius:var(--r);padding:14px}
.btn{border:0;border-radius:12px;padding:10px 14px;font-weight:800;cursor:pointer;background:var(--p);color:#fff}
.btn.ghost{background:#f6f6f6;color:#222}
input[type=file],input[type=text],input[type=number]{width:100%;padding:10px;border:1px solid var(--m);border-radius:12px}
pre{background:#0b0b0c;color:#eaeaea;border-radius:12px;padding:14px;overflow:auto}
small{color:#666}
.grid{display:grid;grid-template-columns:1.1fr .9fr;gap:12px}
@media (max-width:900px){.grid{grid-template-columns:1fr}}
.badge{display:inline-block;padding:2px 8px;border-radius:999px;font-size:12px;font-weight:800}
.ok{background:#ecfdf5;color:#065f46}.mid{background:#fff7ed;color:#92400e}.err{background:#fef2f2;color:#991b1b}
.row{display:grid;grid-template-columns:1fr 1fr;gap:10px}
@media (max-width:720px){.row{grid-template-columns:1fr}}
</style></head><body>
<header><h1>Valutazione automatica · Upload PPTX/PDF</h1></header>
<main class="container">
  <div class="card">
    <form method="post" enctype="multipart/form-data">
      <div class="row">
        <div><label>Carica file (.pptx/.pdf)</label><input type="file" name="file" accept=".pptx,.pdf" required></div>
        <div><label>Executive message atteso (facoltativo)</label><input type="text" name="exec_hint"></div>
      </div>
      <div class="row">
        <div><label>Soglia OK</label><input type="number" name="ok_thr" min="1" max="25" value="{{ok_thr}}"><small>Default 18</small></div>
        <div><label>Soglia Warning</label><input type="number" name="warn_thr" min="1" max="25" value="{{warn_thr}}"><small>Default 14</small></div>
      </div>
      <br><button class="btn" type="submit">Valuta</button>
      <button class="btn ghost" name="action" value="download" type="submit">Scarica ultimo report</button>
    </form>
  </div>
  {% if result %}
  <div class="grid">
    <div class="card">
      <h3>Esito</h3>
      <p><strong>Punteggio:</strong> {{result.total}}/25
      {% if result.total>=ok_thr %}<span class="badge ok">OK ≥{{ok_thr}}</span>
      {% elif result.total>=warn_thr %}<span class="badge mid">Da migliorare</span>
      {% else %}<span class="badge err">Insufficiente</span>{% endif %}</p>
      <ul>
        <li>Chiarezza messaggio: {{result.pts.msg}}/5</li>
        <li>Scelta grafico: {{result.pts.chart}}/5</li>
        <li>Leggibilità: {{result.pts.leg}}/5</li>
        <li>Pulizia visiva: {{result.pts.clean}}/5</li>
        <li>Completezza: {{result.pts.full}}/5</li>
      </ul>
    </div>
    <div class="card"><h3>Feedback</h3><ul>{% for item in result.priorities %}<li>{{item}}</li>{% endfor %}{% if not result.priorities %}<li>Nessuna priorità</li>{% endif %}</ul></div>
  </div>
  <div class="card"><h3>Report</h3><pre>{{result.report}}</pre></div>
  {% endif %}
</main></body></html>
"""

def extract_from_pptx(file_bytes: bytes):
    prs = Presentation(BytesIO(file_bytes))
    texts, has_bar, has_line, has_chart, has_source, has_cta, titles_parlanti = [], False, False, False, False, False, 0
    kpi_counts = []
    for s in prs.slides:
        slide_text = []
        for shp in s.shapes:
            if getattr(shp, "has_text_frame", False):
                t = shp.text_frame.text.strip()
                if t: slide_text.append(t)
            if shp.shape_type == MSO_SHAPE_TYPE.CHART:
                has_chart = True
                try:
                    ctype = shp.chart.chart_type
                    if "COLUMN" in str(ctype) or "BAR" in str(ctype): has_bar = True
                    if "LINE" in str(ctype): has_line = True
                except: pass
        big_text = "\n".join(slide_text)
        texts.append(big_text)
        if re.search(r"\b(aument|dimin|cres|scend|super|sotto|train|resta|sorpass)\w*", big_text, re.I): titles_parlanti += 1
        if re.search(r"\bfonte\b.*(eurostat|dataset|ftth|fttc)", big_text, re.I): has_source = True
        if re.search(r"\b(investire|priorit|agire|ridurre|potenziare|pianificare)\b", big_text, re.I): has_cta = True
        nums = re.findall(r"\b\d+([.,]\d+)?\b", big_text); kpi_counts.append(len(nums))
    return {"text":"\n".join(texts),"titles_parlanti":titles_parlanti,"has_bar":has_bar,"has_line":has_line,
            "has_chart":has_chart,"has_source":has_source,"has_cta":has_cta,"max_nums":max(kpi_counts) if kpi_counts else 0}

def extract_from_pdf(file_bytes: bytes):
    text = extract_text(BytesIO(file_bytes)) or ""
    has_source = bool(re.search(r"\bfonte\b.*(eurostat|dataset|ftth|fttc)", text, re.I))
    has_cta = bool(re.search(r"\b(investire|priorit|agire|ridurre|potenziare|pianificare)\b", text, re.I))
    titles_parlanti = len(re.findall(r"\b(aument|dimin|cres|scend|super|sotto|train|resta|sorpass)\w*", text, re.I))
    kpi_est = len(re.findall(r"\b\d+([.,]\d+)?\b", text))
    return {"text":text,"titles_parlanti":titles_parlanti,"has_bar":False,"has_line":False,"has_chart":False,
            "has_source":has_source,"has_cta":has_cta,"max_nums":kpi_est}

def score(features, exec_hint=""):
    pts = {"msg":0,"chart":0,"leg":0,"clean":0,"full":0}
    c=0
    if features["titles_parlanti"]>=1: c+=3
    if exec_hint and len(exec_hint)>=20: c+=2
    pts["msg"]=min(5,c)
    ch=0
    if features["has_bar"]: ch+=3
    if features["has_line"]: ch+=2
    pts["chart"]=min(5,ch)
    lg=0
    if features["max_nums"]<=12: lg+=3
    if features["titles_parlanti"]>=1: lg+=2
    pts["leg"]=min(5,lg)
    cl=3
    if features["max_nums"]>30: cl-=1
    pts["clean"]=max(0,min(5,cl))
    fu=0
    if features["has_source"]: fu+=3
    if features["has_cta"]: fu+=2
    pts["full"]=min(5,fu)
    total=sum(pts.values())
    priorities=[]
    if pts["msg"]<4: priorities.append("Migliora titoli parlanti e messaggio executive")
    if pts["chart"]<4: priorities.append("Allinea grafici: barre per confronti, linea per trend")
    if pts["leg"]<4: priorities.append("Riduci KPI/slide e usa annotazioni chiare")
    if pts["clean"]<4: priorities.append("Snellisci grafica: niente elementi superflui")
    if pts["full"]<4: priorities.append("Aggiungi fonte e una call-to-action concreta")
    report=f"Totale: {total}/25\n" + "\n".join([f"- {k}: {pts[k]}/5" for k in pts]) + "\nNote automatiche:\n"
    report+=f"- Titoli parlanti: {features['titles_parlanti']}\n- Fonte: {'sì' if features['has_source'] else 'no'}\n- CTA: {'sì' if features['has_cta'] else 'no'}"
    return {"pts":pts,"total":total,"priorities":priorities,"report":report}

@app.route("/",methods=["GET","POST"])
def index():
    ok_thr=int(request.form.get("ok_thr",18) or 18)
    warn_thr=int(request.form.get("warn_thr",14) or 14)
    exec_hint=request.form.get("exec_hint","").strip()
    result=None
    if request.method=="POST":
        action=request.form.get("action")
        if action=="download" and "last_report" in session:
            mem=BytesIO(session["last_report"].encode("utf-8"));mem.seek(0)
            return send_file(mem,mimetype="text/plain",as_attachment=True,download_name="report_valutazione.txt")
        f=request.files.get("file")
        if f:
            name=secure_filename(f.filename);data=f.read()
            if name.lower().endswith(".pptx"): feats=extract_from_pptx(data)
            elif name.lower().endswith(".pdf"): feats=extract_from_pdf(data)
            else: feats={"titles_parlanti":0,"has_bar":False,"has_line":False,"has_chart":False,"has_source":False,"has_cta":False,"max_nums":0}
            s=score(feats,exec_hint=exec_hint);session["last_report"]=s["report"];result=s
    return render_template_string(HTML,result=result,ok_thr=ok_thr,warn_thr=warn_thr)

if __name__=="__main__":
    app.run(host="127.0.0.1",port=5000,debug=False)
